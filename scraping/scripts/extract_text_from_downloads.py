#!/usr/bin/env python3
"""Extract text from files under files/ into extracted_text/<ext>/

Supported types: .txt, .md, .csv, .json, .html, .pdf, .docx, .pptx, .zip
Skips likely-binary or key files (e.g., .key, .pem, images).

Outputs: for each input file creates a UTF-8 .txt file with extracted text.
"""
from __future__ import annotations
import re
import sys
import os
from pathlib import Path
import zipfile
import tempfile
import traceback

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

SKIP_EXT = {'.key', '.pem', '.der', '.p12', '.crt', '.exe', '.dll'}
IMAGE_EXT = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.ico'}

ROOT = Path(__file__).resolve().parents[1]
FILES = ROOT / 'files'
OUTDIR = ROOT / 'extracted_text'

def ensure_out(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    return path

def extract_text_from_pdf(path: Path) -> str:
    if PdfReader is None:
        raise RuntimeError('PyPDF2 not installed')
    text_parts = []
    with path.open('rb') as fh:
        reader = PdfReader(fh)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or '')
            except Exception:
                # continue on extraction errors
                text_parts.append('')
    return '\n\n'.join(text_parts)


def extract_text_from_html(path: Path) -> str:
    # lightweight HTML text extraction without external deps
    data = path.read_text(encoding='utf-8', errors='replace')
    # remove script/style
    data = re.sub(r'(?is)<(script|style).*?>.*?</\1>', '', data)
    # strip tags
    text = re.sub(r'<[^>]+>', '', data)
    return text

def extract_text_from_docx(path: Path) -> str:
    if docx is None:
        raise RuntimeError('python-docx not installed')
    document = docx.Document(path)
    return '\n'.join(p.text for p in document.paragraphs)

def extract_text_from_pptx(path: Path) -> str:
    if Presentation is None:
        raise RuntimeError('python-pptx not installed')
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # many shapes have a text_frame
            try:
                if hasattr(shape, 'text'):
                    txt = shape.text
                elif getattr(shape, 'has_text_frame', False):
                    txt = shape.text_frame.text
                else:
                    txt = ''
            except Exception:
                txt = ''
            if txt:
                parts.append(txt)
    return '\n\n'.join(parts)

def extract_text_from_textfile(path: Path) -> str:
    return path.read_text(encoding='utf-8', errors='replace')

def process_file(path: Path, out_root: Path) -> tuple[bool, str]:
    """Process a single file. Returns (ok, message)."""
    ext = path.suffix.lower()
    rel = path.relative_to(FILES)
    safe_name = str(rel).replace(os.sep, '__')
    out_dir = out_root / (ext.lstrip('.') or 'other')
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / (safe_name + '.txt')

    if ext in SKIP_EXT:
        return False, 'skipped-key-or-binary'
    if ext in IMAGE_EXT:
        return False, 'skipped-image'

    try:
        if ext == '.pdf':
            text = extract_text_from_pdf(path)
        elif ext in {'.htm', '.html'}:
            text = extract_text_from_html(path)
        elif ext in {'.txt', '.md', '.csv', '.json', '.py', '.java', '.c', '.cpp'}:
            text = extract_text_from_textfile(path)
        elif ext == '.docx':
            text = extract_text_from_docx(path)
        elif ext == '.pptx' or ext == '.ppt':
            # try pptx handler; .ppt will be skipped if Presentation cannot read it
            text = extract_text_from_pptx(path)
        elif ext == '.zip':
            # extract to temp and process contents recursively
            with tempfile.TemporaryDirectory() as td:
                with zipfile.ZipFile(path, 'r') as zf:
                    zf.extractall(td)
                # walk extracted
                parts = []
                for root_dir, _, files in os.walk(td):
                    for f in files:
                        p = Path(root_dir) / f
                        try:
                            ok, msg = process_file(p, out_root)
                            if ok:
                                parts.append(f'-- extracted from {f} to {msg}')
                        except Exception:
                            parts.append(f'-- failed {f}: {traceback.format_exc()}')
                text = '\n'.join(parts)
        else:
            return False, f'skipped-unhandled-ext({ext})'

        ensure_out(out_path)
        out_path.write_text(text or '', encoding='utf-8')
        return True, str(out_path)
    except Exception as e:
        return False, f'error:{e}'

def main() -> None:
    if not FILES.exists():
        print('No files/ directory found; nothing to do')
        return
    out_root = OUTDIR
    out_root.mkdir(parents=True, exist_ok=True)

    total = 0
    ok = 0
    failed = 0
    skipped = 0

    for root, _, files in os.walk(FILES):
        for f in files:
            total += 1
            p = Path(root) / f
            success, msg = process_file(p, out_root)
            if success:
                ok += 1
                print('OK:', msg)
            else:
                failed += 1
                print('SKIP/FAIL:', f'{p} -> {msg}')

    print('Done. total=%d ok=%d failed=%d' % (total, ok, failed))

if __name__ == '__main__':
    main()
